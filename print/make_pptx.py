#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
가브리엘라 천사의 집 축복미사 초대장 — A4 세로 1장 PPTX 생성기

모든 텍스트를 개별 텍스트 상자로 배치해 누구나 클릭해서 고칠 수 있게 만든다.
수정 후 재생성:  python3 make_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT  = os.path.join(HERE, "가브리엘라_축복미사_초대장_A4.pptx")

# ── 색 ────────────────────────────────────────────────
INK       = RGBColor(0x1C, 0x2A, 0x25)
INK_SOFT  = RGBColor(0x57, 0x66, 0x5F)
INK_FAINT = RGBColor(0x8A, 0x96, 0x8F)
GREEN     = RGBColor(0x15, 0x7A, 0x63)
LINE      = RGBColor(0xDD, 0xD5, 0xC4)

# ── 글꼴 ──────────────────────────────────────────────
SERIF = "나눔명조"      # 표제·본문
SANS  = "나눔고딕"      # 라벨·안내

# ── 지면 ──────────────────────────────────────────────
PW, PH   = 210.0, 297.0
MARGIN_X = 20.0
CW       = PW - MARGIN_X * 2       # 본문 폭 170mm

prs = Presentation()
prs.slide_width  = Mm(PW)
prs.slide_height = Mm(PH)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # 빈 레이아웃


def set_font(run, name, size_pt, bold=False, color=INK, spacing=None):
    """한글이 제대로 나오도록 latin/ea/cs 글꼴을 모두 지정한다."""
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)
    if spacing is not None:                 # 자간 (1/100 pt)
        rPr.set("spc", str(int(spacing * 100)))


def textbox(x, y, w, h, lines, align=PP_ALIGN.CENTER, line_spacing=None,
            space_after=0):
    """lines = [(텍스트, 글꼴, 크기, 굵기, 색, 자간), ...] — 한 항목이 한 문단"""
    tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, spec in enumerate(lines):
        text, font, size, bold, color = spec[:5]
        spacing = spec[5] if len(spec) > 5 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        set_font(run, font, size, bold, color, spacing)
    return tb


def rule(y, weight_pt=0.5, color=LINE, x=MARGIN_X, w=CW):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Mm(x), Mm(y),
                                    Mm(x + w), Mm(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight_pt)
    return ln


# ══════════════════════════════════════════════════════
# 1. 로고
# ══════════════════════════════════════════════════════
LOGO_W = 18.0
LOGO_H = LOGO_W * 166 / 215
slide.shapes.add_picture(os.path.join(HERE, "logo.png"),
                         Mm((PW - LOGO_W) / 2), Mm(12), Mm(LOGO_W), Mm(LOGO_H))

# ══════════════════════════════════════════════════════
# 2. 머리말 · 표제
# ══════════════════════════════════════════════════════
textbox(MARGIN_X, 29.5, CW, 6,
        [("말기암환자 무료쉼터", SANS, 8.5, True, GREEN, 1.6)])

textbox(MARGIN_X, 36.5, CW, 12,
        [("가브리엘라 천사의 집", SERIF, 23, True, INK)])

textbox(MARGIN_X, 50.5, CW, 8,
        [("축복미사", SERIF, 12.5, False, INK_SOFT)])

# ══════════════════════════════════════════════════════
# 3. 일시 · 장소
# ══════════════════════════════════════════════════════
rule(61.5, 1.1, INK)

textbox(MARGIN_X, 65, CW, 9,
        [("2026년 9월 17일 목요일 오전 11시", SERIF, 15, True, INK)])

textbox(MARGIN_X, 75.5, CW, 6,
        [("가브리엘라 천사의 집", SERIF, 10, True, INK)])

textbox(MARGIN_X, 81.5, CW, 6,
        [("인천광역시 강화군 화도면 안골길61번길 17-12 (내리 904)",
          SERIF, 10, False, INK_SOFT)])

rule(89.5)

# ══════════════════════════════════════════════════════
# 4. 조감도
# ══════════════════════════════════════════════════════
PIC_Y, PIC_H = 93.0, 30.0
pic = slide.shapes.add_picture(os.path.join(HERE, "hero_crop.jpg"),
                               Mm(MARGIN_X), Mm(PIC_Y), Mm(CW), Mm(PIC_H))
pic.line.color.rgb = LINE
pic.line.width = Pt(0.4)

textbox(MARGIN_X, PIC_Y + PIC_H + 2, CW, 5,
        [("가브리엘라 천사의 집 조감도", SERIF, 7.5, False, INK_FAINT)])

# ══════════════════════════════════════════════════════
# 5. 초대의 글
# ══════════════════════════════════════════════════════
LY = 140.0

textbox(MARGIN_X, LY, CW, 14,
        [("환자와 가족이 삶의 마지막 여정을 두려움이 아닌", SERIF, 10.5, True, INK),
         ("사랑과 감사 속에서 준비할 수 있도록 돕는 무료 돌봄시설입니다.", SERIF, 10.5, True, INK)],
        line_spacing=1.5)

textbox(MARGIN_X, LY + 16, CW, 8,
        [("죽음을 기다리는 곳이 아니라 삶을 완성하는 곳!", SERIF, 11, True, GREEN)],
        line_spacing=1.5)

textbox(MARGIN_X, LY + 26, CW, 20,
        [("종교와 관계없이 무료 쉼터와 돌봄을 제공하는 이곳은", SERIF, 10, False, INK),
         ("소중한 분들의 기도와 성원에 힘입어", SERIF, 10, False, INK),
         ("하느님 사랑의 터전을 이루게 되었습니다.", SERIF, 10, False, INK)],
        line_spacing=1.5)

textbox(MARGIN_X, LY + 46, CW, 8,
        [("이제 첫 문을 여는 자리에 소중한 여러분을 모시고자 합니다.", SERIF, 10, False, INK)],
        line_spacing=1.5)

# ══════════════════════════════════════════════════════
# 6. 식순
# ══════════════════════════════════════════════════════
PY = 198.0
rule(PY)

textbox(MARGIN_X, PY + 4, CW, 6,
        [("식  순", SANS, 8, True, GREEN, 2.2)])

PROG = [
    ("11:00", "축복미사",
     "인천교구장 정신철 요한 세례자 주교 주례 · 염수정 안드레아 추기경, 사제단 공동집전"),
    ("12:20", "기념식", "재단 임·직원, 내·외빈 참석"),
    ("이후",  "오찬",   "식사가 제공됩니다"),
]
ROW_Y, ROW_H = PY + 13, 8.0
for i, (t, name, desc) in enumerate(PROG):
    y = ROW_Y + i * ROW_H
    textbox(MARGIN_X, y, 16, 6, [(t, SERIF, 10, True, GREEN)], align=PP_ALIGN.LEFT)
    tb = slide.shapes.add_textbox(Mm(MARGIN_X + 21), Mm(y), Mm(CW - 21), Mm(6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = name + "   "
    set_font(r1, SERIF, 10, True, INK)
    r2 = p.add_run(); r2.text = desc
    set_font(r2, SERIF, 9, False, INK_SOFT)
    if i < len(PROG) - 1:
        rule(y + ROW_H - 2, 0.3)

# ══════════════════════════════════════════════════════
# 7. 안내 (2단)
# ══════════════════════════════════════════════════════
NY = ROW_Y + len(PROG) * ROW_H + 4
rule(NY)

COL_W = (CW - 6) / 2
# 각 항목은 줄 단위로 끊어 둔다 — 단어가 중간에 잘리지 않도록
NOTICE = [
    ("준비해 주세요",
     ["야외미사 집전입니다. 양산·모자 등 햇빛을 가릴 것과 개인",
      "텀블러를 지참해 주시면 좋겠습니다."],
     None),
    ("주차 안내",
     ["가브리엘라 천사의 집 안에는 주차 공간이 부족합니다.",
      "마니산 공영주차장에 주차해 주시면, 행사장까지 이동 차량을",
      "운행합니다."],
     "인천광역시 강화군 화도면 마니산로675번길 18 (상방리 405-4)"),
]
for i, (title, body, addr) in enumerate(NOTICE):
    x = MARGIN_X + i * (COL_W + 6)
    textbox(x, NY + 4, COL_W, 5,
            [(title, SANS, 7.5, True, GREEN, 1.8)], align=PP_ALIGN.LEFT)
    textbox(x, NY + 9.5, COL_W, 14,
            [(ln, SERIF, 9, False, INK) for ln in body],
            align=PP_ALIGN.LEFT, line_spacing=1.45)
    if addr:
        textbox(x, NY + 24, COL_W, 6,
                [(addr, SERIF, 8, False, INK_SOFT)], align=PP_ALIGN.LEFT, line_spacing=1.4)

# ══════════════════════════════════════════════════════
# 8. 꼬리말
# ══════════════════════════════════════════════════════
FY = NY + 32
rule(FY, 1.1, INK)
textbox(MARGIN_X, FY + 3.5, CW, 7,
        [("재단법인 마뗄암재단", SERIF, 9.5, True, INK)])
textbox(MARGIN_X, FY + 10, CW, 5,
        [("Mater Cancer Foundation · www.mcancer.com", SANS, 7.5, False, INK_FAINT, 0.6)])

prs.save(OUT)
print("생성:", OUT)
print("마지막 요소 하단: %.1fmm / 지면 %.0fmm" % (FY + 15.5, PH))
